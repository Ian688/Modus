"""Office document tools (Phase A3): Excel / Word / PPT for everyday users.

The user's common workload is analysing a multi-thousand-row Excel sheet,
tweaking Word formatting, or building a PPT.  The existing text tools cannot
read these binary formats (read_file hard-refuses >1MB and is text-only).  These
handlers parse the binary in-process and return bounded summaries/edits to the
model, with the full result persisted as an artifact.

Security posture:
- Read tools: ``data_disclosure="workspace_content"`` (cell/paragraph text goes
  to the model), ``capabilities=("filesystem",)``, safe + auto-ALLOW.
- Write tools: ``is_read_only=False`` + ``danger_level="medium"`` so they go
  through the approval gate.
- Bounded by construction: Excel reads are streaming (``read_only=True`` +
  ``iter_rows``), row/column caps everywhere, model payloads are small summaries
  while the full sheet is persisted as an artifact.
- The 1MB read_file limit is intentionally bypassed: Excel/Word/PPT are ZIP
  binaries and were never readable as text; these handlers parse them directly.
"""

from __future__ import annotations

import os
from pathlib import Path
from typing import Any

from modus.tools.base import ToolContext, ToolResult

_EXCEL_MAX_ROWS = 50_000
_EXCEL_MAX_COLS = 200
_PREVIEW_ROWS = 10
_PREVIEW_COLS = 12
_ARTIFACT_CHARS = 20_000


def _resolve_ws_path(context: ToolContext, value: str) -> Path:
    """Resolve a workspace path via the PathGuard (home-anchored boundary)."""
    from modus.tools.builtins import _resolve_path

    return _resolve_path(context, value)


def _persist(context: ToolContext, kind: str, title: str, content: str) -> dict[str, Any] | None:
    from modus.tools.builtins import _persist_tool_result

    return _persist_tool_result(context, kind=kind, title=title, content=content)


def _bounded(text: str, limit: int = 8000) -> str:
    if len(text) <= limit:
        return text
    return text[:limit] + f"\n… [{len(text) - limit} chars omitted]"


# ── Excel ──


async def excel_analyze(payload: dict[str, Any], context: ToolContext) -> ToolResult:
    """Analyze an Excel workbook: sheets, dimensions, header, sample rows, stats."""
    path_value = str(payload.get("path") or "").strip()
    if not path_value:
        return ToolResult("excel_analyze requires a path", is_error=True)
    try:
        path = _resolve_ws_path(context, path_value)
    except Exception as exc:
        return ToolResult(f"excel_analyze path error: {exc}", is_error=True)
    if not path.exists():
        return ToolResult(f"file not found: {path_value}", is_error=True)
    try:
        import openpyxl

        wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
    except Exception as exc:
        return ToolResult(f"excel_analyze failed to open: {exc}", is_error=True)

    lines: list[str] = [f"Workbook: {path_value}", f"Sheets: {wb.sheetnames}"]
    try:
        ws = wb.active
        lines.append(f"Active sheet: {ws.title} ({ws.max_row} rows × {ws.max_column} cols)")
        # Header row + preview rows
        header: list[str] = []
        preview_rows: list[list[str]] = []
        for i, row in enumerate(ws.iter_rows(min_row=1, max_row=_PREVIEW_ROWS,
                                            max_col=min(ws.max_column or 1, _PREVIEW_COLS),
                                            values_only=True)):
            cells = ["" if v is None else str(v) for v in row]
            if i == 0:
                header = cells
            else:
                preview_rows.append(cells)
        lines.append(f"Header: {header}")
        lines.append("Preview rows:")
        for r in preview_rows:
            lines.append("  " + " | ".join(r))
        # Simple numeric stats for the first few numeric columns
        lines.append(_excel_numeric_stats(ws, header))
        wb.close()
    except Exception as exc:
        lines.append(f"(preview error: {exc})")
    text = _bounded("\n".join(lines))
    artifact = _persist(context, "office", f"excel · {path.name}", text)
    return ToolResult(
        text,
        display_summary=f"分析 Excel：{path_value}",
        metadata={"operation": "excel_analyze", "path": str(path), "sheets": len(wb.sheetnames)},
        artifacts=[artifact] if artifact else [],
    )


def _excel_numeric_stats(ws, header: list[str]) -> str:
    """Mean/min/max for up to 3 numeric columns, sampled across the sheet."""
    if not header:
        return ""
    numeric_cols: list[tuple[int, str]] = []
    for idx, name in enumerate(header):
        if idx >= 5:
            break
        numeric_cols.append((idx + 1, str(name)))
    stats_lines: list[str] = []
    for col_idx, col_name in numeric_cols:
        values: list[float] = []
        try:
            for row in ws.iter_rows(min_row=2, max_col=col_idx, max_row=_EXCEL_MAX_ROWS,
                                    values_only=True):
                v = row[-1] if row else None
                if isinstance(v, (int, float)) and not isinstance(v, bool):
                    values.append(float(v))
        except Exception:
            continue
        if len(values) >= 2:
            stats_lines.append(
                f"  {col_name}: n={len(values)} mean={sum(values)/len(values):.2f} "
                f"min={min(values):.2f} max={max(values):.2f}"
            )
    return ("Numeric stats:\n" + "\n".join(stats_lines)) if stats_lines else ""


async def excel_query(payload: dict[str, Any], context: ToolContext) -> ToolResult:
    """Filter rows of an Excel sheet by a column value / numeric range."""
    path_value = str(payload.get("path") or "").strip()
    if not path_value:
        return ToolResult("excel_query requires a path", is_error=True)
    column = str(payload.get("column") or "").strip()
    equals = payload.get("equals")
    gt = payload.get("gt")
    limit = max(1, min(int(payload.get("limit") or 20), 200))
    if not column or (equals is None and gt is None):
        return ToolResult("excel_query requires column + (equals or gt)", is_error=True)
    try:
        path = _resolve_ws_path(context, path_value)
    except Exception as exc:
        return ToolResult(f"excel_query path error: {exc}", is_error=True)
    try:
        import openpyxl

        wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
        ws = wb.active
    except Exception as exc:
        return ToolResult(f"excel_query failed: {exc}", is_error=True)

    # Find the header index matching the column name.
    header_row = next(ws.iter_rows(min_row=1, max_row=1, values_only=True), ())
    col_idx = None
    for idx, name in enumerate(header_row):
        if str(name or "").strip().lower() == column.lower():
            col_idx = idx
            break
    if col_idx is None:
        wb.close()
        return ToolResult(f"column not found: {column} (header: {header_row})", is_error=True)

    rows: list[str] = []
    matched = 0
    for row in ws.iter_rows(min_row=2, values_only=True):
        if col_idx >= len(row):
            continue
        v = row[col_idx]
        try:
            if equals is not None:
                if str(v or "") == str(equals):
                    hit = True
                else:
                    hit = False
            else:
                hit = isinstance(v, (int, float)) and float(v) > float(gt)
        except (TypeError, ValueError):
            hit = False
        if hit:
            rows.append(" | ".join("" if c is None else str(c) for c in row))
            matched += 1
            if matched >= limit:
                rows.append(f"... [limited to {limit} rows]")
                break
    wb.close()
    text = _bounded("\n".join(rows) or f"(no rows match {column}={equals or '>' + str(gt)})")
    return ToolResult(
        text,
        display_summary=f"查询 Excel：{path_value}",
        metadata={"operation": "excel_query", "path": str(path), "matched": matched},
    )


# ── Word ──


async def word_extract(payload: dict[str, Any], context: ToolContext) -> ToolResult:
    """Extract paragraphs, headings, and table cell text from a .docx."""
    path_value = str(payload.get("path") or "").strip()
    if not path_value:
        return ToolResult("word_extract requires a path", is_error=True)
    try:
        path = _resolve_ws_path(context, path_value)
    except Exception as exc:
        return ToolResult(f"word_extract path error: {exc}", is_error=True)
    try:
        import docx

        d = docx.Document(str(path))
    except Exception as exc:
        return ToolResult(f"word_extract failed: {exc}", is_error=True)
    lines: list[str] = []
    for p in d.paragraphs[:100]:
        t = p.text.strip()
        if t:
            style = p.style.name if p.style else ""
            prefix = "# " if "Heading" in style else ""
            lines.append(prefix + t)
    for i, table in enumerate(d.tables[:5]):
        lines.append(f"[Table {i+1}]")
        for row in table.rows[:5]:
            lines.append("  " + " | ".join(cell.text.strip()[:40] for cell in row.cells))
    text = _bounded("\n".join(lines) or "(empty document)")
    artifact = _persist(context, "office", f"word · {path.name}", text)
    return ToolResult(
        text,
        display_summary=f"提取 Word：{path_value}",
        metadata={"operation": "word_extract", "path": str(path)},
        artifacts=[artifact] if artifact else [],
    )


async def word_edit(payload: dict[str, Any], context: ToolContext) -> ToolResult:
    """Edit a .docx: replace an exact text occurrence across paragraphs."""
    path_value = str(payload.get("path") or "").strip()
    find = str(payload.get("find") or "").strip()
    replace = str(payload.get("replace") or "")
    if not path_value or not find:
        return ToolResult("word_edit requires path + find", is_error=True)
    try:
        path = _resolve_ws_path(context, path_value)
    except Exception as exc:
        return ToolResult(f"word_edit path error: {exc}", is_error=True)
    try:
        import docx

        d = docx.Document(str(path))
    except Exception as exc:
        return ToolResult(f"word_edit failed: {exc}", is_error=True)
    replaced = 0
    for p in d.paragraphs:
        if find in p.text:
            # Replace across runs to preserve formatting as much as possible.
            for run in p.runs:
                if find in run.text:
                    run.text = run.text.replace(find, replace)
                    replaced += 1
    if replaced == 0:
        return ToolResult(f"text not found: {find}", is_error=True)
    try:
        d.save(str(path))
    except Exception as exc:
        return ToolResult(f"word_edit save failed: {exc}", is_error=True)
    return ToolResult(
        f"Replaced {replaced} occurrence(s) of '{find}' in {path_value}",
        display_summary=f"编辑 Word：{path_value}",
        metadata={"operation": "word_edit", "path": str(path), "replaced": replaced},
    )


# ── PPT ──


async def pptx_extract(payload: dict[str, Any], context: ToolContext) -> ToolResult:
    """Extract slide titles and text bodies from a .pptx."""
    path_value = str(payload.get("path") or "").strip()
    if not path_value:
        return ToolResult("pptx_extract requires a path", is_error=True)
    try:
        path = _resolve_ws_path(context, path_value)
    except Exception as exc:
        return ToolResult(f"pptx_extract path error: {exc}", is_error=True)
    try:
        from pptx import Presentation

        prs = Presentation(str(path))
    except Exception as exc:
        return ToolResult(f"pptx_extract failed: {exc}", is_error=True)
    lines: list[str] = []
    for i, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                t = shape.text_frame.text.strip()
                if t:
                    texts.append(t)
        if texts:
            lines.append(f"Slide {i}: " + " | ".join(texts[:5]))
    text = _bounded("\n".join(lines) or "(empty presentation)")
    artifact = _persist(context, "office", f"ppt · {path.name}", text)
    return ToolResult(
        text,
        display_summary=f"提取 PPT：{path_value}",
        metadata={"operation": "pptx_extract", "path": str(path), "slides": len(prs.slides)},
        artifacts=[artifact] if artifact else [],
    )


async def pptx_build(payload: dict[str, Any], context: ToolContext) -> ToolResult:
    """Build a simple .pptx from a list of title/body slide specs."""
    path_value = str(payload.get("path") or "").strip()
    slides = payload.get("slides")
    if not path_value or not isinstance(slides, list) or not slides:
        return ToolResult("pptx_build requires path + slides (list of {title, body?})", is_error=True)
    try:
        path = _resolve_ws_path(context, path_value)
    except Exception as exc:
        return ToolResult(f"pptx_build path error: {exc}", is_error=True)
    try:
        from pptx import Presentation
        from pptx.util import Inches

        prs = Presentation()
        for spec in slides[:20]:
            if not isinstance(spec, dict):
                continue
            slide = prs.slides.add_slide(prs.slide_layouts[1])  # title+content
            title = str(spec.get("title") or "")[:120]
            body = str(spec.get("body") or "")[:1000]
            if slide.shapes.title:
                slide.shapes.title.text = title
            if body and slide.placeholders:
                slide.placeholders[1].text = body
        prs.save(str(path))
    except Exception as exc:
        return ToolResult(f"pptx_build failed: {exc}", is_error=True)
    return ToolResult(
        f"Built {len(slides[:20])} slides into {path_value}",
        display_summary=f"生成 PPT：{path_value}",
        metadata={"operation": "pptx_build", "path": str(path), "slides": len(slides[:20])},
    )
